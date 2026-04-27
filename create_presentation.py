"""
Bank Leumi Employee Rights & Options - PPTX Presentation Generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy
from lxml import etree

# ── Brand colors ──────────────────────────────────────────────────────────────
LEUMI_BLUE      = RGBColor(0x00, 0x4B, 0x9B)
LEUMI_DARK_BLUE = RGBColor(0x00, 0x28, 0x64)
LEUMI_LIGHT     = RGBColor(0xD6, 0xE8, 0xFF)
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT       = RGBColor(0x33, 0x33, 0x44)

LOGO_PATH = "/home/user/text/leumi_logo.png"
SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)

# ── Slide content ─────────────────────────────────────────────────────────────
SLIDES = [
    # 0 – Title slide
    {
        "type": "title",
        "title": "זכויות ואפשרויות עובדי בנק לאומי",
        "subtitle": "מדריך מקיף לעובד הלאומי | 2024",
    },
    # 1 – Agenda
    {
        "type": "agenda",
        "title": "נושאי המצגת",
        "items": [
            "תנאי העסקה ושכר",
            "ימי חופשה ומחלה",
            "פנסיה וקרן השתלמות",
            "זכויות הורות",
            "הטבות רווחה ובריאות",
            "פיתוח מקצועי והשכלה",
            "זכויות ועד עובדים",
            "סיום יחסי עבודה",
        ],
    },
    # 2
    {
        "type": "content",
        "title": "תנאי העסקה ושכר",
        "icon": "💰",
        "points": [
            "שכר מינימום ענפי גבוה משמעותית מהמינימום החוקי",
            "הסכם קיבוצי מיוחד לעובדי הבנקים (הסכם ענפי)",
            "תוספות שכר: ותק, תפקיד, תפוקה ובונוסים שנתיים",
            "דירוג ומסלול קידום שקוף ומוגדר",
            "גמול עבור שעות נוספות לפי חוק עבודת שעות נוספות",
            "תוספת מחיה ותוספות בגין תפקידי ניהול",
            "שכר 13 – מענק שנתי מוסכם בהסכם הקיבוצי",
        ],
    },
    # 3
    {
        "type": "content",
        "title": "ימי חופשה ומחלה",
        "icon": "🏖️",
        "points": [
            "חופשה שנתית: 14–28 ימי עסקים (בהתאם לוותק)",
            "ימי מחלה: 1.5 יום לחודש ₪ צבירה ללא הגבלה",
            "ימי אבל (שבעה): עד 7 ימים בגין קרוב משפחה מדרגה ראשונה",
            "ימי אבל בגין קרוב משפחה מדרגה שנייה: עד 3 ימים",
            "מחלת ילד: עד 8 ימים בשנה (חפיפה עם ימי מחלה)",
            "ימי בחירה (ימים אישיים): 2 ימים נוספים בשנה",
            "חגים: כל ימי החג היהודי + חגי מיעוטים לפי הצורך",
        ],
    },
    # 4
    {
        "type": "content",
        "title": "פנסיה וקרן השתלמות",
        "icon": "🏦",
        "points": [
            "פנסיה מקיפה: הפרשות מעסיק 7.5% + פיצויים 8.33%",
            "הפרשות עובד לפנסיה: 6% משכר הברוטו",
            "קרן השתלמות: מעסיק 7.5% + עובד 2.5% (פטורה ממס לאחר 6 שנים)",
            "ביטוח מנהלים לעובדים בתפקידי ניהול בכיר",
            "זכות לנייד קרן פנסיה קודמת ללא קנס",
            "ייעוץ פנסיוני חינם דרך הבנק",
            "קרן פנסיה ענפית עם תשואה מובטחת (עד גיל 50)",
        ],
    },
    # 5
    {
        "type": "content",
        "title": "זכויות הורות",
        "icon": "👶",
        "points": [
            "חופשת לידה: 26 שבועות לאם (15 שבועות בתשלום מביטוח לאומי)",
            "חופשת לידה לאב: 7 ימי עבודה בתשלום מלא",
            "שמירת היריון: בתשלום מלא על פי הוראות הרופא",
            "חזרה לעבודה: הגנה מפיטורין 60 יום לאחר החזרה",
            "הקלות לאם מניקה: הפסקות הנקה בשכר",
            "קצובת לידה מהבנק בנוסף על קצבת ביטוח לאומי",
            "גמישות בשעות עבודה לאחר חזרה מחופשת לידה",
        ],
    },
    # 6
    {
        "type": "content",
        "title": "הטבות רווחה ובריאות",
        "icon": "❤️",
        "points": [
            "ביטוח בריאות קבוצתי מסובסד לעובד ומשפחתו",
            "ביטוח חיים ואובדן כושר עבודה על חשבון הבנק",
            "רכב חברה או קצובת רכב לתפקידי ניהול",
            "הלוואות עובדים בריבית מועדפת",
            "מנוי לחדר כושר / פיטנס בסבסוד",
            "סיוע בשכר לימוד לילדים (מלגות)",
            "ארועי חברה ופעילויות גיבוש שנתיות",
            "מתנות לחגים וסל חג",
        ],
    },
    # 7
    {
        "type": "content",
        "title": "פיתוח מקצועי והשכלה",
        "icon": "📚",
        "points": [
            "תוכניות הכשרה פנימיות: הבנק משקיע עשרות שעות הדרכה לעובד",
            "לימודים אקדמיים: מימון חלקי לתואר ראשון ושני",
            "תוכניות Fast Track לקידום מהיר של עובדים מצטיינים",
            "קורסים מקצועיים ורישיונות: ביטוח, השקעות, פיננסים",
            "מנטורינג ואימון אישי לעובדים בכירים",
            "קרן לימודים בקרן ההשתלמות ניתנת למשיכה לצרכי השכלה",
            "כנסים מקצועיים בארץ ובחו\"ל",
        ],
    },
    # 8
    {
        "type": "content",
        "title": "זכויות ועד עובדים",
        "icon": "🤝",
        "points": [
            "הסכם קיבוצי ענפי ומפעלי – כוח מיקוח חזק",
            "נציגות ועד פעיל בכל סניף ויחידה",
            "זכות שביתה מוגנת על פי חוק",
            "הגנה מפני פיטורין שרירותיים",
            "טיפול בתלונות ואפליה במקום העבודה",
            "ייצוג משפטי חינם בסכסוכי עבודה",
            "ועדת שימוע חובה לפני כל פיטורין",
        ],
    },
    # 9
    {
        "type": "content",
        "title": "סיום יחסי עבודה",
        "icon": "📋",
        "points": [
            "פיצויי פיטורין: שכר חודשי לכל שנת וותק",
            "הודעה מוקדמת: 1 יום לחודש עד שנה, לאחר מכן חודש",
            "תוכנית פרישה מוקדמת עם תנאים מיוחדים (גיל 57+)",
            "גמול פרישה נוסף מעבר לחוק – לפי ההסכם הקיבוצי",
            "שחרור קרן פנסיה וקרן השתלמות",
            "המשך ביטוח בריאות למשך 6 חודשים לאחר עזיבה",
            "שמירה על אי תחרות: 12 חודש בתחום הבנקאות (ניתן לערעור)",
        ],
    },
    # 10 – Summary
    {
        "type": "summary",
        "title": "סיכום – כוחו של עובד לאומי",
        "bullets": [
            ("📞", "ועד העובדים", "פנה/י לנציג הוועד בסניף שלך"),
            ("⚖️", "ייעוץ משפטי", "זכאי/ת לייצוג ללא עלות בסכסוכי עבודה"),
            ("📄", "ההסכם הקיבוצי", "בקש/י עותק מהמנהל או מהוועד"),
            ("💡", "הכרת הזכויות", "ידע הוא כוח – קרא/י את ההסכם הקיבוצי"),
        ],
    },
    # 11 – Thank you
    {
        "type": "thankyou",
        "title": "תודה!",
        "subtitle": "עובד מכיר זכויותיו – עובד מוגן",
        "footer": "בנק לאומי לישראל בע\"מ | 2024",
    },
]

# ── Helper functions ──────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.RIGHT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_logo(slide, size=Inches(1.1)):
    try:
        slide.shapes.add_picture(LOGO_PATH,
                                 SLIDE_W - size - Inches(0.2),
                                 Inches(0.15),
                                 size, size)
    except Exception:
        pass


def blue_bar(slide, top, height, color=LEUMI_BLUE):
    """Horizontal colored rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), top, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bullet_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, WHITE)

    # Top bar
    blue_bar(slide, Inches(0), Inches(1.2), LEUMI_BLUE)

    # Title
    add_text_box(slide, data["title"],
                 Inches(0.4), Inches(0.1),
                 Inches(10), Inches(1.0),
                 font_size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT)

    # Icon circle accent
    icon_box = slide.shapes.add_shape(1, Inches(11.6), Inches(0.1),
                                      Inches(1.5), Inches(1.0))
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = LEUMI_DARK_BLUE
    icon_box.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(11.6), Inches(0.1),
                                   Inches(1.5), Inches(1.0))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = data.get("icon", "")
    r.font.size = Pt(36)

    # Bottom bar
    blue_bar(slide, SLIDE_H - Inches(0.45), Inches(0.45), LEUMI_DARK_BLUE)
    add_text_box(slide, "בנק לאומי | זכויות ואפשרויות עובדים",
                 Inches(0.3), SLIDE_H - Inches(0.4),
                 Inches(12), Inches(0.38),
                 font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)

    # Bullet points
    points = data.get("points", [])
    top_start = Inches(1.35)
    row_h = (SLIDE_H - top_start - Inches(0.55)) / max(len(points), 1)

    for i, point in enumerate(points):
        top = top_start + i * row_h

        # Bullet dot
        dot = slide.shapes.add_shape(1,
                                     Inches(12.5), top + row_h * 0.3,
                                     Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = LEUMI_BLUE
        dot.line.fill.background()

        # Alternating light-blue row
        if i % 2 == 0:
            row_bg = slide.shapes.add_shape(1,
                                            Inches(0.3), top + Inches(0.05),
                                            Inches(12.7), row_h - Inches(0.08))
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = LEUMI_LIGHT
            row_bg.line.fill.background()

        add_text_box(slide, point,
                     Inches(0.5), top,
                     Inches(12.1), row_h,
                     font_size=16, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)

    add_logo(slide)
    return slide


# ── Build slides ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

for data in SLIDES:
    stype = data["type"]
    blank = prs.slide_layouts[6]

    # ── Title slide ──────────────────────────────────────────────────────────
    if stype == "title":
        slide = prs.slides.add_slide(blank)
        set_bg(slide, LEUMI_BLUE)

        # Dark blue bottom half
        blue_bar(slide, SLIDE_H / 2, SLIDE_H / 2, LEUMI_DARK_BLUE)

        # Decorative white wave strip
        wave = slide.shapes.add_shape(1, Inches(0),
                                      SLIDE_H / 2 - Inches(0.15),
                                      SLIDE_W, Inches(0.3))
        wave.fill.solid()
        wave.fill.fore_color.rgb = WHITE
        wave.line.fill.background()

        add_logo(slide, size=Inches(2.2))

        add_text_box(slide, data["title"],
                     Inches(0.5), Inches(1.6),
                     Inches(11), Inches(2.0),
                     font_size=44, bold=True, color=WHITE,
                     align=PP_ALIGN.RIGHT)

        add_text_box(slide, data["subtitle"],
                     Inches(0.5), SLIDE_H / 2 + Inches(0.4),
                     Inches(11), Inches(1.2),
                     font_size=24, color=LEUMI_LIGHT,
                     align=PP_ALIGN.RIGHT, italic=True)

    # ── Agenda slide ─────────────────────────────────────────────────────────
    elif stype == "agenda":
        slide = prs.slides.add_slide(blank)
        set_bg(slide, WHITE)

        blue_bar(slide, Inches(0), Inches(1.3), LEUMI_DARK_BLUE)
        add_text_box(slide, data["title"],
                     Inches(0.4), Inches(0.15),
                     Inches(12), Inches(1.0),
                     font_size=34, bold=True, color=WHITE,
                     align=PP_ALIGN.RIGHT)

        # Left accent bar
        accent = slide.shapes.add_shape(1,
                                        Inches(0), Inches(1.3),
                                        Inches(0.25), SLIDE_H - Inches(1.3))
        accent.fill.solid()
        accent.fill.fore_color.rgb = LEUMI_BLUE
        accent.line.fill.background()

        items = data["items"]
        col_items = [items[:4], items[4:]]
        col_lefts = [Inches(7.0), Inches(0.5)]

        for col_idx, col in enumerate(col_items):
            for row_idx, item in enumerate(col):
                top = Inches(1.6) + row_idx * Inches(1.3)
                left = col_lefts[col_idx]

                # Number badge
                badge = slide.shapes.add_shape(1,
                                               left + Inches(4.8),
                                               top + Inches(0.1),
                                               Inches(0.5), Inches(0.5))
                badge.fill.solid()
                badge.fill.fore_color.rgb = LEUMI_BLUE
                badge.line.fill.background()

                num_tbx = slide.shapes.add_textbox(left + Inches(4.8),
                                                   top + Inches(0.05),
                                                   Inches(0.5), Inches(0.5))
                nt = num_tbx.text_frame
                np_ = nt.paragraphs[0]
                np_.alignment = PP_ALIGN.CENTER
                nr = np_.add_run()
                nr.text = str(col_idx * 4 + row_idx + 1)
                nr.font.size = Pt(14)
                nr.font.bold = True
                nr.font.color.rgb = WHITE

                add_text_box(slide, item, left, top,
                             Inches(4.7), Inches(0.7),
                             font_size=18, color=GRAY_TEXT,
                             align=PP_ALIGN.RIGHT)

        blue_bar(slide, SLIDE_H - Inches(0.45), Inches(0.45), LEUMI_DARK_BLUE)
        add_text_box(slide, "בנק לאומי | זכויות ואפשרויות עובדים",
                     Inches(0.3), SLIDE_H - Inches(0.4),
                     Inches(12), Inches(0.38),
                     font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)
        add_logo(slide)

    # ── Content slides ───────────────────────────────────────────────────────
    elif stype == "content":
        bullet_slide(prs, data)

    # ── Summary slide ────────────────────────────────────────────────────────
    elif stype == "summary":
        slide = prs.slides.add_slide(blank)
        set_bg(slide, LEUMI_DARK_BLUE)

        blue_bar(slide, Inches(0), Inches(1.3), LEUMI_BLUE)
        add_text_box(slide, data["title"],
                     Inches(0.4), Inches(0.15),
                     Inches(12), Inches(1.0),
                     font_size=34, bold=True, color=WHITE,
                     align=PP_ALIGN.RIGHT)

        for i, (icon, heading, body) in enumerate(data["bullets"]):
            top = Inches(1.6) + i * Inches(1.3)

            card = slide.shapes.add_shape(1,
                                          Inches(0.4), top,
                                          Inches(12.5), Inches(1.15))
            card.fill.solid()
            card.fill.fore_color.rgb = LEUMI_BLUE
            card.line.fill.background()

            add_text_box(slide, f"{icon}  {heading}: {body}",
                         Inches(0.6), top + Inches(0.1),
                         Inches(12), Inches(0.95),
                         font_size=20, bold=False, color=WHITE,
                         align=PP_ALIGN.RIGHT)

        add_logo(slide)

    # ── Thank-you slide ──────────────────────────────────────────────────────
    elif stype == "thankyou":
        slide = prs.slides.add_slide(blank)
        set_bg(slide, LEUMI_BLUE)

        blue_bar(slide, SLIDE_H / 2, SLIDE_H / 2, LEUMI_DARK_BLUE)

        wave = slide.shapes.add_shape(1, Inches(0),
                                      SLIDE_H / 2 - Inches(0.1),
                                      SLIDE_W, Inches(0.2))
        wave.fill.solid()
        wave.fill.fore_color.rgb = WHITE
        wave.line.fill.background()

        add_logo(slide, size=Inches(2.0))

        add_text_box(slide, data["title"],
                     Inches(0.5), Inches(2.2),
                     Inches(12), Inches(1.5),
                     font_size=60, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, data["subtitle"],
                     Inches(0.5), SLIDE_H / 2 + Inches(0.5),
                     Inches(12), Inches(1.0),
                     font_size=26, color=LEUMI_LIGHT,
                     align=PP_ALIGN.CENTER, italic=True)

        add_text_box(slide, data["footer"],
                     Inches(0.5), SLIDE_H - Inches(0.7),
                     Inches(12), Inches(0.5),
                     font_size=13, color=WHITE,
                     align=PP_ALIGN.CENTER)

OUTPUT = "/home/user/text/leumi_employee_rights.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}  ({len(SLIDES)} slides)")
